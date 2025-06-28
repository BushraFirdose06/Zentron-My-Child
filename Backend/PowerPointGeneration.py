from pptx import Presentation
from pptx.util import Inches
from dotenv import dotenv_values
import os

env_vars = dotenv_values(".env")
Username = env_vars.get("Username")
Assistantname = env_vars.get("Assistantname")

def GeneratePresentation(topic, content=None, slides=5):
    """Generates a PowerPoint (.pptx) file from a given topic."""
    try:
        prs = Presentation()  # Create a new PowerPoint

        # === Slide 1: Title Slide ===
        title_slide = prs.slide_layouts[0]  # 0 = Title Slide Layout
        slide = prs.slides.add_slide(title_slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = topic
        subtitle.text = f"Created by {Assistantname} for {Username}"

        # === If no content is provided, generate it using ChatBot ===
        if not content:
            from Backend.Chatbot import ChatBot
            content = ChatBot(f"Generate detailed content for a PowerPoint presentation about {topic}")

        # === Split content into slides ===
        paragraphs = [p for p in content.split('\n') if p.strip()]
        slides_content = []
        for i in range(0, len(paragraphs), slides):
            slides_content.append("\n".join(paragraphs[i:i+slides]))

        # === Add Content Slides ===
        for text in slides_content:
            content_slide = prs.slide_layouts[1]  # 1 = Title and Content
            slide = prs.slides.add_slide(content_slide)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = topic
            tf = body.text_frame
            for line in text.split('\n'):
                p = tf.add_paragraph()
                p.text = line
                p.level = 0  # Bullet level

        # === Save the PowerPoint ===
        filename = f"Data/{topic.replace(' ', '_')}_presentation.pptx"
        prs.save(filename)
        return filename

    except Exception as e:
        print(f"❌ Error in PowerPoint Generation: {e}")
        return None